"""PPTX text extraction utilities for Knowledge Base indexing.

This module provides functionality to extract text content from PowerPoint files
for indexing into AWS Bedrock Knowledge Base, which doesn't natively support PPTX format.
"""

from typing import List
from pptx import Presentation


def extract_text_from_pptx(file_path: str) -> str:
    """Extract all text content from a PowerPoint file.

    Extracts text from all slides including:
    - Slide titles
    - Body text and bullet points
    - Text in shapes and text boxes
    - Notes (speaker notes)

    Args:
        file_path: Path to the PPTX file

    Returns:
        str: Extracted text content formatted as markdown-style sections

    Raises:
        FileNotFoundError: If file doesn't exist
        ValueError: If file is not a valid PPTX
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"Invalid PowerPoint file: {str(e)}") from e

    text_parts: List[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text: List[str] = [f"# Slide {slide_num}"]

        # Extract text from all shapes on the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        # Extract speaker notes if present
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_text.append(f"\n**Notes:** {notes_text}")

        # Only add slide if it has content
        if len(slide_text) > 1:
            text_parts.append("\n\n".join(slide_text))

    return "\n\n---\n\n".join(text_parts)
